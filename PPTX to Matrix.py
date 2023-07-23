from pptx import Presentation
import pandas as pd
import numpy as np
import os

#Folder that contains each participant's files
ppt_folder = r'C:\Example\Slides'

#File that containes the mapping between the letters and audio links
label_path = r'C:\Example\Letters to Links.xlsx'

#Matrix Output Location 
output_path = r'C:\Example\Matrix_Output.xlsx'


#The slide numbers start at 0 i.e. slides 4 & 5 need to be entered as 3 & 4 respectivly
slide_no_1 = 4
slide_no_2 = 3

#The pair threshold is how close the boxes can be to each other to be considered pairs.
# A value of 1 indicated that two boxes are only pairs if their top left corners are
# within 1x of the height or width of each other (to the left/right or above/below).
#A value of 1.25 would be within 1.25x of the height or width so there can be a small gap.

pair_threshold = 0.25

labels = pd.read_excel(label_path)
letter_list = list(labels['Letters'])
speaker_list = list(labels['Speaker'])
link_list = list(labels['Audio Link'])

label_dict = dict(zip(link_list,tuple(zip(speaker_list,letter_list))))



def matrix_generator(slide,label_dict,pair_threshold):
# This function contains the logic to creat a matrix

    location_list = []
    name_list = []

    #Each shape that is in the "Picture" format has its location and dimension recorded. The audio link name is also recorded. 
    for shape in slide.shapes:

        
        if shape.__class__.__name__ == 'Picture':

            address = shape.click_action.hyperlink.address

            name = label_dict[address]
            
            name_list.append(name)
            location_list.append([shape.left,shape.top,shape.width,shape.height])

        else:

            next

    # Converts name_list to multi-level
    name_list = pd.MultiIndex.from_tuples(name_list)

    count = len(location_list)

    pair_array = np.empty([count,count])

    
    #Pairs are generated for every rectangle that is next to another rectangle. This is determined by the x and y-axis locations. 
    for column in range(0,count):

        for row in range(0,count):

            x_loc =   location_list[row][0]
            x_0_min = location_list[column][0]-(location_list[column][2]*pair_threshold)
            x_0_max = location_list[column][0]+(location_list[column][2]*pair_threshold)
            x_min =   location_list[column][0]-(location_list[column][2]*(1+pair_threshold))
            x_max =   location_list[column][0]+(location_list[column][2]*(1+pair_threshold))                                  

            y_loc =   location_list[row][1]
            y_0_min = location_list[column][1]-(location_list[column][3]*pair_threshold)
            y_0_max = location_list[column][1]+(location_list[column][3]*pair_threshold)
            y_min =   location_list[column][1]-(location_list[column][3]*(1+pair_threshold))
            y_max =   location_list[column][1]+(location_list[column][3]*(1+pair_threshold))

            y_match = (x_loc < x_0_max) & (x_loc > x_0_min) & (y_loc < y_max) & (y_loc > y_min)

            x_match = (y_loc < y_0_max) & (y_loc > y_0_min) & (x_loc < x_max) & (x_loc > x_min)
                                                        
                      
            if x_match | y_match:

                pair_array[row,column] = 1

            else:

                pair_array[row,column] = 0


    set_list = []
    group_list = []

    #Each pair in the pair array is checked agains other pairs to capture common pairs and generate groups.
    for a in range(0,len(pair_array)):

        set_a = set(list(np.where(pair_array[a,:]==1)[0]))

        for b in range(0,len(pair_array)):

            set_b = set(list(np.where(pair_array[b,:]==1)[0]))

            if not(set_a.isdisjoint(set_b)):

                set_a = (set_a | set_b) 
        
        set_list.append(set_a)
        group_list.append(list(set_a))


    group_list = np.unique(group_list)


    group_change = True
    change_count = 0

    #The groups are checked against each other until there are no more common pairs found. This code is likely redundant in most cases. 
    while group_change:

        for a in range(0,len(group_list)):
            
            set_main = set(group_list[a])

            for b in range(0,len(group_list)):
                
                set_comp = set(group_list[b]) 

                if not(set_main.isdisjoint(set_comp)):
                
                    set_main = (set_main | set_comp)

                    change_count+1

            group_list[a] = list(set_main)


            if change_count == 0 and a == (len(group_list) -1):

                group_change = False

    group_list = np.unique(group_list)

    group_array = pair_array.copy()

    for group in group_list:

        for i in range(0,len(pair_array)):

            if i in group:
                
                for member in group:

                    if i == member:

                        group_array[i][member] = 0

                    else: 

                        group_array[i][member] = 1


    group_matrix = pd.DataFrame(group_array, columns = name_list, index = name_list)

    return group_matrix


def matrix_to_excel(folder,output_path,label_dict,slide_no_1,slide_no_2,pair_threshold):
#This function formats the group matricies, creates an excel file, and pastes the results from each powerpoint. 

    file_no = 0

    writer = pd.ExcelWriter(output_path,engine='xlsxwriter')
    workbook=writer.book
    worksheet=workbook.add_worksheet('Result')
    writer.sheets['Result'] = worksheet

    row = 0
    

    for file in os.listdir(folder):

        if file[-5:] != '.pptx':

            continue
    
        try: 
            input_path = folder+'\\'+file
            prs = Presentation(input_path)

        except:

            input_path = folder+'/'+file
            prs = Presentation(input_path)


        slide1 = prs.slides[slide_no_1]
        slide2 = prs.slides[slide_no_2]

        matrix1 = matrix_generator(slide1,label_dict,pair_threshold)
        matrix2 = matrix_generator(slide2,label_dict,pair_threshold)


        column = matrix1.shape[1]
        
        worksheet.write_string(row+2, 0, file[:-5])

        worksheet.write_string(row+3, 0, matrix1.columns[0][1][0])
        matrix1.to_excel(writer,sheet_name='Result',startrow=(row+3) , startcol=0)

        worksheet.write_string(row+3, column + 3, matrix2.columns[0][1][0])
        matrix2.to_excel(writer,sheet_name='Result',startrow=(row+3), startcol=column + 3)

        file_no = file_no + 1
        row = row+column+5

    writer.save()
    writer.close()    
    return 


matrix_to_excel(ppt_folder,output_path,label_dict,slide_no_1,slide_no_2,pair_threshold)

print('Done')
